#!/usr/bin/env python3
"""
AUTOSALES Modernization - Project Completion Deck
Generates a professional PPTX summarizing the completed modernization.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)
MEDIUM_BLUE = RGBColor(0x1B, 0x4F, 0x72)
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED = RGBColor(0xF4, 0x43, 0x36)
TEAL = RGBColor(0x00, 0x96, 0x88)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, text, top=0, height=Inches(1.0)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_subtitle_bar(slide, text, top=Inches(1.0), height=Inches(0.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MEDIUM_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_kpi_card(slide, label, value, left, top, width=Inches(2.5), height=Inches(1.5), color=ACCENT_BLUE):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    # Value
    add_text_box(slide, value, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.7),
                 font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, label, left + Inches(0.2), top + Inches(0.9), width - Inches(0.4), Inches(0.4),
                 font_size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)


def add_table(slide, data, left, top, width, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

            if r == 0:  # Header
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = MEDIUM_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY

    return table


def create_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ═══════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    add_text_box(slide, "AUTOSALES", Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
                 font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Mainframe Modernization — Project Completion Report",
                 Inches(0.5), Inches(2.7), Inches(12), Inches(0.8),
                 font_size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "IMS DC / COBOL / DB2  →  Spring Boot + React + PostgreSQL",
                 Inches(0.5), Inches(3.5), Inches(12), Inches(0.6),
                 font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, "Project #4  |  279 Legacy Files  |  ~95,200 SLOC  |  10 Business Domains",
                 Inches(0.5), Inches(4.6), Inches(12), Inches(0.5),
                 font_size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "March 2026",
                 Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
                 font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════
    # SLIDE 2: Key Metrics Dashboard
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Project Completion — Key Metrics")
    add_subtitle_bar(slide, "All 103+ legacy programs modernized with 12 beyond-parity enhancements")

    # Row 1: Main metrics
    y1 = Inches(1.8)
    add_kpi_card(slide, "Legacy Programs Ported", "103+", Inches(0.5), y1, color=GREEN)
    add_kpi_card(slide, "Java Source Files", "364", Inches(3.3), y1, color=ACCENT_BLUE)
    add_kpi_card(slide, "React Pages", "50", Inches(6.1), y1, color=PURPLE)
    add_kpi_card(slide, "REST Endpoints", "187", Inches(8.9), y1, color=TEAL)
    add_kpi_card(slide, "Unit Tests", "374", Inches(11.0), y1, width=Inches(1.8), color=GREEN)

    # Row 2: Code metrics
    y2 = Inches(3.6)
    add_kpi_card(slide, "Java SLOC", "19,000", Inches(0.5), y2, color=ACCENT_BLUE)
    add_kpi_card(slide, "TypeScript SLOC", "21,000", Inches(3.3), y2, color=PURPLE)
    add_kpi_card(slide, "Legacy Total SLOC", "95,200", Inches(6.1), y2, color=ORANGE)
    add_kpi_card(slide, "Code Reduction", "58%", Inches(8.9), y2, color=GREEN)
    add_kpi_card(slide, "Enhancements", "12/12", Inches(11.0), y2, width=Inches(1.8), color=TEAL)

    # Row 3: Component breakdown
    y3 = Inches(5.4)
    add_kpi_card(slide, "JPA Entities", "51", Inches(0.5), y3, width=Inches(2.0), color=MEDIUM_BLUE)
    add_kpi_card(slide, "Repositories", "53", Inches(2.8), y3, width=Inches(2.0), color=MEDIUM_BLUE)
    add_kpi_card(slide, "Services", "38", Inches(5.1), y3, width=Inches(2.0), color=MEDIUM_BLUE)
    add_kpi_card(slide, "Controllers", "32", Inches(7.4), y3, width=Inches(2.0), color=MEDIUM_BLUE)
    add_kpi_card(slide, "DTOs", "124", Inches(9.7), y3, width=Inches(2.0), color=MEDIUM_BLUE)
    add_kpi_card(slide, "Migrations", "28", Inches(12.0), y3, width=Inches(1.0), color=MEDIUM_BLUE)

    # ═══════════════════════════════════════════════════════════
    # SLIDE 3: Wave Execution Summary
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Wave Execution Summary")
    add_subtitle_bar(slide, "8 waves completed — progressive build from foundation to full functionality")

    data = [
        ["Wave", "Scope", "Legacy Programs", "Services", "Endpoints", "Tests", "React Pages"],
        ["0: Foundation", "Auth, Schema, Utilities", "16 common libs", "—", "2", "44", "3"],
        ["1: Admin", "Dealers, Pricing, Tax, Models", "7", "7", "30", "73", "7"],
        ["2: Customer", "CRUD, Credit, Leads", "7", "3", "12", "87", "4"],
        ["3: Sales", "Deal Lifecycle State Machine", "8", "1 (550 LOC)", "9", "99", "3"],
        ["4: Finance", "Finance, F&I, Floor Plan", "12", "5", "15", "168", "8"],
        ["5: Vehicle", "Inventory, Stock, Production", "26", "6", "55", "230", "14"],
        ["6: Registration", "Registration, Warranty, Recall", "11", "4", "30", "374", "7"],
        ["7: Batch", "Daily/Monthly/Weekly, Integrations", "11", "9+", "37", "374", "2"],
        ["TOTAL", "", "103+", "38", "187", "374", "50"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(1.5), Inches(3.0), Inches(1.5), Inches(1.5), Inches(1.2), Inches(1.0), Inches(1.2)])

    # ═══════════════════════════════════════════════════════════
    # SLIDE 4: Architecture
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Architecture: Modular Monolith")
    add_subtitle_bar(slide, "Why NOT microservices: deal lifecycle spans 6+ modules atomically")

    # Architecture decision box
    add_text_box(slide, "Key Decision: Modular Monolith over Microservices",
                 Inches(0.5), Inches(1.8), Inches(6), Inches(0.5),
                 font_size=18, bold=True, color=DARK_BLUE)

    reasons = [
        "• Deal completion atomically touches Customer + Vehicle + Stock + Finance + Registration + Warranty",
        "• Single @Transactional boundary — no saga patterns, no distributed locks",
        "• Modules cleanly separated by package — extractable to microservices later",
        "• Docker deployment: 3 containers (app + db + ui) — minutes, not hours",
    ]
    for i, r in enumerate(reasons):
        add_text_box(slide, r, Inches(0.7), Inches(2.5 + i * 0.4), Inches(6), Inches(0.4),
                     font_size=13, color=DARK_GRAY)

    # Technology stack
    add_text_box(slide, "Technology Stack",
                 Inches(7.5), Inches(1.8), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color=DARK_BLUE)

    stack = [
        ["Layer", "Technology"],
        ["Backend", "Spring Boot 3.3 / Java 21"],
        ["Database", "PostgreSQL 16 (Flyway)"],
        ["Security", "Spring Security + JWT"],
        ["Frontend", "React 18 + TypeScript 5"],
        ["Styling", "Tailwind CSS 3"],
        ["Build", "Maven + Vite"],
        ["Deploy", "Docker Compose"],
        ["Testing", "JUnit 5 + Mockito"],
    ]
    add_table(slide, stack, Inches(7.5), Inches(2.4), Inches(5.3))

    # Module diagram
    add_text_box(slide, "8 Domain Modules",
                 Inches(0.5), Inches(4.5), Inches(6), Inches(0.5),
                 font_size=16, bold=True, color=DARK_BLUE)

    modules = [
        ("Admin", GREEN), ("Customer", ACCENT_BLUE), ("Sales", ORANGE), ("Finance", PURPLE),
        ("Vehicle", TEAL), ("Registration", MEDIUM_BLUE), ("Batch", MED_GRAY), ("Floor Plan", GREEN),
    ]
    for i, (name, color) in enumerate(modules):
        col = i % 4
        row = i // 4
        left = Inches(0.7 + col * 1.5)
        top = Inches(5.2 + row * 0.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.3), Inches(0.5))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True

    # ═══════════════════════════════════════════════════════════
    # SLIDE 5: Before & After
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Before & After: IMS DC Terminal → Modern Web Application")
    add_subtitle_bar(slide, "16 MFS green-screen terminals replaced by 50 responsive React pages")

    data = [
        ["Aspect", "IMS DC (Before)", "React + Spring Boot (After)"],
        ["Interface", "Fixed 80×24 green text terminal", "Responsive web, any screen size"],
        ["Navigation", "Transaction codes (VEHI, FPLR)", "Collapsible sidebar with 8 groups"],
        ["Data Entry", "Tab between fixed fields", "Modern forms with real-time validation"],
        ["Output", "Fixed-width monospace text", "Color badges, charts, cards, tables"],
        ["Pagination", "PF7/PF8 keys, 12 rows max", "Click pagination, configurable page size"],
        ["Export", "Print screen or SYSOUT", "CSV download, one click"],
        ["Access", "3270 terminal emulator only", "Any web browser, any device"],
        ["Deployment", "JCL / IEBCOPY / SMP/E", "Docker Compose (3 containers, 2 minutes)"],
        ["Testing", "Manual only", "374 automated unit tests"],
        ["User Roles", "IMS sign-on", "JWT + RBAC (5 roles: Admin/Mgr/Sales/Finance/Clerk)"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(1.8), Inches(5.0), Inches(5.5)])

    # ═══════════════════════════════════════════════════════════
    # SLIDE 6: CLIST/REXX Modernization
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "CLIST & REXX Automation Scripts — Architectural Elimination")
    add_subtitle_bar(slide, "13 mainframe scripts (2,300 SLOC) replaced by platform capabilities — not ported line-by-line")

    data = [
        ["Legacy Script", "Type", "Purpose", "Modern Equivalent"],
        ["ASSUBMIT", "REXX", "Submit batch JCL with date override", "BatchJobsPage 'Run Now' buttons"],
        ["ASSTATUS", "REXX", "Check job status via SDSF", "Dashboard System Health panel"],
        ["ASBROWSE", "CLIST", "Browse report output in ISPF", "BatchReportsPage (6 tabs)"],
        ["ASVINLKP", "REXX", "VIN lookup with deal history", "VehicleDetailPage + VinDecodePanel"],
        ["ASDEALER", "REXX", "Dealer KPI dashboard", "StockDashboardPage (7 KPI cards)"],
        ["ASMENU", "CLIST", "ISPF main menu panel", "Collapsible Sidebar (8 nav groups)"],
        ["ASLOGON", "CLIST", "TSO logon, allocate datasets", "docker compose up (2 min)"],
        ["ASCOMPIL", "CLIST", "DB2 precompile + COBOL + link", "mvn compile (one command)"],
        ["ASEDIT", "CLIST", "ISPF editor (cols 7-72)", "VS Code / IntelliJ IDE"],
        ["ASMIGRTE", "REXX", "DEV->QA->PROD (IEBCOPY)", "Docker image tag + push"],
        ["ASBACKUP", "REXX", "DB2 UNLOAD + DFDSS backup", "pg_dump / Docker volumes"],
        ["ASDBCHK", "REXX", "DB2 health check", "Spring Actuator /health endpoint"],
        ["ASIMSCHK", "REXX", "IMS queue depth monitoring", "Dashboard health badges"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(1.3), Inches(0.8), Inches(4.0), Inches(4.5)])

    add_text_box(slide, "Total Legacy: 95,200 SLOC  (COBOL 74.3K + MFS 8.3K + JCL 4.4K + Copybooks 3.9K + REXX 1.5K + DDL 1.1K + CLIST 0.8K + PSB/DBD 0.9K)",
                 Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 font_size=12, bold=True, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════
    # SLIDE 7: Quality & Defect Metrics
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Quality & Defect Metrics")
    add_subtitle_bar(slide, "19 defects found and fixed — zero in core business logic translation")

    y1 = Inches(1.8)
    add_kpi_card(slide, "Total Tests", "374", Inches(0.5), y1, color=GREEN)
    add_kpi_card(slide, "Pass Rate", "100%", Inches(3.3), y1, color=GREEN)
    add_kpi_card(slide, "Defects Found", "19", Inches(6.1), y1, color=ORANGE)
    add_kpi_card(slide, "Biz Logic Defects", "0", Inches(8.9), y1, color=GREEN)
    add_kpi_card(slide, "Build Cycles", "30+", Inches(11.0), y1, width=Inches(1.8), color=ACCENT_BLUE)

    data = [
        ["Category", "Count", "Critical", "High", "Medium", "Low", "Examples"],
        ["Deployment/Infra", "7", "1", "4", "1", "0", "PG reserved word, Docker volume, auth token"],
        ["Code (non-logic)", "10", "0", "1", "4", "5", "Seed overflow, @Auditable on read TX, FormField"],
        ["UX/UI", "2", "0", "0", "2", "0", "Sidebar length, Dashboard hardcoded data"],
        ["Business Logic", "0", "0", "0", "0", "0", "Zero defects in COBOL->Java translation"],
        ["TOTAL", "19", "1", "5", "7", "5", ""],
    ]
    add_table(slide, data, Inches(0.3), Inches(3.6), Inches(12.7),
              [Inches(2.0), Inches(0.8), Inches(0.9), Inches(0.8), Inches(1.0), Inches(0.7), Inches(5.5)])

    add_text_box(slide, "Key: All defects were in deployment config, seed data, or component wiring. "
                 "The COBOL-to-Java business logic (deal lifecycle, financial calcs, stock management) translated with zero defects.",
                 Inches(0.5), Inches(6.3), Inches(12), Inches(0.6),
                 font_size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════
    # SLIDE 8: 12 Enhancements Beyond Parity
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "12 Enhancements Beyond Mainframe Parity")
    add_subtitle_bar(slide, "Modernization as opportunity, not just migration")

    data = [
        ["#", "Enhancement", "Business Value", "Status"],
        ["1", "Sales Pipeline Dashboard", "Visual deal flow — managers see bottlenecks", "✓ Done"],
        ["2", "Inventory Health Dashboard", "Aging heat map — prevents aged inventory", "✓ Done"],
        ["3", "Finance Calculator Suite", "Side-by-side term comparison for customers", "✓ Done"],
        ["4", "Floor Plan Exposure Dashboard", "Lender breakdown + curtailment alerts", "✓ Done"],
        ["5", "Audit Trail UI", "Searchable audit log — compliance ready", "✓ Done"],
        ["6", "Excel/CSV Export", "One-click data export from any list view", "✓ Done"],
        ["7", "Real-time Deal Profitability", "Live gross calculation during negotiation", "✓ Done"],
        ["8", "VIN Instant Decode", "Auto-decode manufacturer/year/plant on entry", "✓ Done"],
        ["9", "On-demand Reports", "Run any batch report instantly from UI", "✓ Done"],
        ["10", "System Health Dashboard", "Batch job monitoring with health badges", "✓ Done"],
        ["11", "User Management UI", "Full RBAC admin — create, lock, reset password", "✓ Done"],
        ["12", "Paginated + Searchable Lists", "Server-side pagination on all 50 pages", "✓ Done"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(0.5), Inches(2.8), Inches(5.5), Inches(1.0)])

    # ═══════════════════════════════════════════════════════════
    # SLIDE 7: Lessons Learned
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Key Lessons & Technical Decisions")
    add_subtitle_bar(slide, "First IMS DC modernization — different from prior CICS-based projects")

    data = [
        ["Category", "Lesson", "Impact"],
        ["Architecture", "Modular monolith > microservices for cross-domain transactions", "Zero distributed TX issues"],
        ["IMS DC", "MFS screens → React pages + REST endpoints (1:1 mapping)", "Clean porting pattern"],
        ["Data", "5 IMS hierarchical DBDs flattened to relational tables", "All data consolidated"],
        ["Financial Math", "BigDecimal mandatory — COBOL DEC(11,2) precision preserved", "30+ calc tests pass"],
        ["State Machine", "11 vehicle status codes with strict transition matrix", "No invalid transitions"],
        ["Docker/Windows", "Volume mount hot-reload fails on Windows — bake source into image", "Reliable dev workflow"],
        ["PostgreSQL", "Reserved words (system_user), type mismatches need vigilance", "DDL/JPA alignment"],
        ["Parallel Build", "Wave 6 + 7 built simultaneously — zero conflicts with file ownership rules", "2x throughput"],
        ["AI-Assisted", "Claude Code for parallel agents, legacy doc analysis, test generation", "Accelerated delivery"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(1.5), Inches(6.5), Inches(4.0)])

    # ═══════════════════════════════════════════════════════════
    # SLIDE 8: Risk Mitigation
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Risk Mitigation Results")
    add_subtitle_bar(slide, "All identified risks successfully mitigated")

    data = [
        ["Risk", "Severity", "Mitigation", "Result"],
        ["IMS hierarchical data flattening", "HIGH", "Mapped every DBD segment chain", "5 DBDs flattened ✓"],
        ["Financial calc precision", "HIGH", "BigDecimal with exact COBOL algorithms", "30+ tests pass ✓"],
        ["Deal state machine edge cases", "HIGH", "Comprehensive transition matrix", "12 tests, all passing ✓"],
        ["Cross-module TX consistency", "MEDIUM", "Single @Transactional (monolith)", "Zero issues ✓"],
        ["Floor plan day-count variance", "MEDIUM", "Support all 3 bases (360/365/ACT)", "Calculator tested ✓"],
        ["Multi-jurisdiction tax", "MEDIUM", "State/county/city rate engine", "Tax tests pass ✓"],
        ["RAM constraint (dev machine)", "MEDIUM", "Monolith: single JVM ~512MB", "Docker runs fine ✓"],
    ]
    add_table(slide, data, Inches(0.3), Inches(1.7), Inches(12.7),
              [Inches(3.0), Inches(1.2), Inches(4.5), Inches(3.0)])

    # ═══════════════════════════════════════════════════════════
    # SLIDE 9: What This Demonstrates
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "What This Project Demonstrates")
    add_subtitle_bar(slide, "Capabilities validated for future engagements")

    # Left column: Practice capabilities
    add_text_box(slide, "For the Modernization Practice",
                 Inches(0.5), Inches(1.8), Inches(6), Inches(0.5),
                 font_size=18, bold=True, color=DARK_BLUE)

    practice = [
        "✓ IMS DC expertise — first IMS-based modernization (prior 3 were CICS)",
        "✓ Dual data store consolidation — IMS hierarchical + DB2 → single PostgreSQL",
        "✓ Architecture flexibility — chose monolith over microservices based on requirements",
        "✓ AI-accelerated delivery — parallel development with Claude Code agents",
        "✓ Beyond-parity value — 12 enhancements that never existed on mainframe",
    ]
    for i, p in enumerate(practice):
        add_text_box(slide, p, Inches(0.7), Inches(2.5 + i * 0.5), Inches(6), Inches(0.4),
                     font_size=13, color=DARK_GRAY)

    # Right column: Stakeholder value
    add_text_box(slide, "For Stakeholders",
                 Inches(7.5), Inches(1.8), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color=DARK_BLUE)

    stakeholder = [
        "✓ 100% functional parity — every mainframe TX has modern equivalent",
        "✓ 50 polished pages vs 16 green-screen terminals",
        "✓ Self-service: on-demand reports, CSV export, dashboards",
        "✓ Operational visibility: batch monitoring, audit trail",
        "✓ Reduced risk: 374 automated tests, containerized deployment",
        "✓ Future-ready: API-first, extractable modules, standard stack",
    ]
    for i, s in enumerate(stakeholder):
        add_text_box(slide, s, Inches(7.7), Inches(2.5 + i * 0.5), Inches(5), Inches(0.4),
                     font_size=13, color=DARK_GRAY)

    # ═══════════════════════════════════════════════════════════
    # SLIDE 10: Thank You
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    add_text_box(slide, "AUTOSALES Modernization Complete",
                 Inches(0.5), Inches(2.0), Inches(12), Inches(1.0),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "103+ COBOL Programs  •  374 Tests  •  50 React Pages  •  12 Enhancements",
                 Inches(0.5), Inches(3.2), Inches(12), Inches(0.6),
                 font_size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "IMS DC / COBOL / DB2  →  Spring Boot 3.3 + React 18 + PostgreSQL 16",
                 Inches(0.5), Inches(4.0), Inches(12), Inches(0.6),
                 font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.8), Inches(5), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, "Project #4  |  March 2026",
                 Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
                 font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), "AUTOSALES_Project_Completion_Deck.pptx")
    prs.save(output_path)
    print(f"Deck saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_deck()
