#!/usr/bin/env python3
"""
AUTOSALES Modernization - Executive Summary Approval Deck
Generates a professional PPTX for stakeholder approval.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)      # Title bars
MEDIUM_BLUE = RGBColor(0x1B, 0x4F, 0x72)     # Subtitle bars, headers
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)     # Highlights, links
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)      # Light backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED = RGBColor(0xF4, 0x43, 0x36)
TEAL = RGBColor(0x00, 0x96, 0x88)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_bar(slide, text, top=0, height=Inches(1.0)):
    """Add a dark blue title bar across the top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    # Indent
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return bar


def add_subtitle_bar(slide, text, top=Inches(1.0), height=Inches(0.5)):
    """Add a medium blue subtitle bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MEDIUM_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return bar


def add_text_box(slide, text, left, top, width, height, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_stat_box(slide, label, value, left, top, width=Inches(2.5),
                 height=Inches(1.5), color=ACCENT_BLUE):
    """Add a stat box with big number and label."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Value (big)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = color

    # Label (small)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(12)
    run2.font.color.rgb = MED_GRAY

    return box


def add_card(slide, title, items, left, top, width=Inches(3.8),
             height=Inches(3.5), accent=ACCENT_BLUE):
    """Add a card with title and bullet items."""
    # Card background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.width = Pt(1)

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    # Title
    add_text_box(slide, title, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4), font_size=16, color=accent, bold=True)

    # Items
    for i, item in enumerate(items):
        add_text_box(slide, f"• {item}", left + Inches(0.2),
                     top + Inches(0.6) + Inches(i * 0.35),
                     width - Inches(0.4), Inches(0.35), font_size=11, color=DARK_GRAY)

    return box


def add_table_slide(slide, headers, rows, left, top, col_widths):
    """Add a table to slide."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top,
        sum(col_widths), Inches(0.4) * (len(rows) + 1)
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

    return table_shape


# ═══════════════════════════════════════════════════════════════
# SLIDE GENERATORS
# ═══════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(2.8), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, "MAINFRAME MODERNIZATION", Inches(1), Inches(1.5),
                 Inches(11), Inches(0.6), font_size=18, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, "AUTOSALES — Automotive Dealer\nSales & Reporting Platform",
                 Inches(1), Inches(3.0), Inches(11), Inches(1.5),
                 font_size=36, color=WHITE, bold=True)

    add_text_box(slide, "Executive Summary & Modernization Approval",
                 Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                 font_size=20, color=RGBColor(0xBB, 0xBB, 0xBB))

    add_text_box(slide, "IMS DC  •  COBOL  •  DB2  •  IMS Hierarchical DB  •  REXX  •  CLIST  •  MFS  •  EDI",
                 Inches(1), Inches(5.5), Inches(11), Inches(0.4),
                 font_size=14, color=MED_GRAY)

    add_text_box(slide, "Project #4 in AI-Assisted Modernization Series  |  March 2026",
                 Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 font_size=12, color=MED_GRAY)


def slide_02_exec_summary(prs):
    """Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Executive Summary")
    add_subtitle_bar(slide, "The largest and most complex mainframe modernization in the series")

    # Stat boxes row
    stats = [
        ("COBOL Programs", "103+", ACCENT_BLUE),
        ("Total SLOC", "~84,500", RED),
        ("Business Domains", "10", TEAL),
        ("DB2 Tables", "46+", GREEN),
        ("IMS Databases", "5", ORANGE),
        ("External Integrations", "7", PURPLE),
    ]
    for i, (label, value, color) in enumerate(stats):
        add_stat_box(slide, label, value,
                     Inches(0.2) + Inches(i * 2.2), Inches(1.8),
                     width=Inches(2.0), height=Inches(1.3), color=color)

    # Key message
    add_text_box(slide, "AUTOSALES is a comprehensive automotive dealer sales platform running on "
                 "IBM mainframe with IMS DC transaction processing, dual data stores (DB2 + IMS hierarchical), "
                 "and deep integration with external systems (EDI, CRM, DMS, GL, Data Lake).",
                 Inches(0.5), Inches(3.4), Inches(12), Inches(0.8),
                 font_size=14, color=DARK_GRAY)

    # Why modernize boxes
    reasons = [
        ("Aging Platform Risk", "IMS DC skills shrinking\nVendor support declining\nHardware costs rising"),
        ("Business Agility", "3270 green screens limit UX\nBatch-only reporting\nNo mobile/web access"),
        ("Integration Friction", "EDI-only interfaces\nFixed-format file exchanges\nNo real-time APIs"),
        ("Competitive Pressure", "Modern dealer platforms\nCloud-native competitors\nCustomer expectations"),
    ]
    for i, (title, desc) in enumerate(reasons):
        left = Inches(0.4) + Inches(i * 3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.5),
                                      Inches(2.9), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()
        add_text_box(slide, title, left + Inches(0.15), Inches(4.6),
                     Inches(2.6), Inches(0.4), font_size=14, color=MEDIUM_BLUE, bold=True)
        add_text_box(slide, desc, left + Inches(0.15), Inches(5.1),
                     Inches(2.6), Inches(1.5), font_size=11, color=MED_GRAY)


def slide_03_current_state(prs):
    """Current State - System Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Current State — System Overview")
    add_subtitle_bar(slide, "IMS DC / COBOL / DB2 / IMS Hierarchical DB / REXX / CLIST / MFS / EDI")

    # Technology stack table
    headers = ["Component", "Technology", "Count", "SLOC"]
    rows = [
        ["Online Programs", "IMS DC / COBOL", "76", "48,474"],
        ["Batch Programs", "IMS BMP / COBOL / DB2", "11", "7,003"],
        ["Common Modules", "COBOL subroutines", "16", "8,904"],
        ["Copybooks", "DCLGEN + Common", "43+", "3,912"],
        ["Screen Definitions", "MFS", "16", "8,274"],
        ["Job Control", "JCL", "11", "4,443"],
        ["Relational Database", "DB2 z/OS", "46+ tables", "—"],
        ["Hierarchical Database", "IMS (DBD/PSB)", "5 DBDs / 16 PSBs", "898"],
        ["Automation Scripts", "REXX + CLIST", "13", "2,601"],
        ["EDI Integration", "ANSI X12 (214/856)", "2 formats", "—"],
    ]
    col_widths = [Inches(2.5), Inches(3), Inches(1.5), Inches(3.5)]
    add_table_slide(slide, headers, rows, Inches(0.8), Inches(1.8), col_widths)


def slide_04_business_domains(prs):
    """Business Domain Map."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Business Domain Map")
    add_subtitle_bar(slide, "10 domains covering the complete automotive dealer sales lifecycle")

    domains = [
        ("Customer\nManagement", "7 programs", ["Add/Edit/Search", "Credit Pre-Qual", "Lead Tracking", "Purchase History"], ACCENT_BLUE),
        ("Sales\nProcess", "8 programs", ["Quote/Worksheet", "Negotiation", "Trade-in Eval", "Approval Workflow", "Complete/Cancel"], GREEN),
        ("Finance &\nFloor Plan", "12 programs", ["Loan Calculator", "Lease Calculator", "F&I Products", "Floor Plan Mgmt", "Credit Check"], ORANGE),
        ("Vehicle &\nInventory", "24+ programs", ["Production→Ship→Receive", "Stock Position", "Aging Analysis", "Transfer/PDI", "EDI Tracking"], TEAL),
        ("Registration", "5 programs", ["Generate/Validate", "DMV Submission", "Status Tracking"], PURPLE),
        ("Warranty &\nRecall", "6 programs", ["Auto-Registration", "Recall Campaigns", "Notifications", "Claims Report"], RED),
        ("Administration", "8 programs", ["Dealer Master", "Tax Rates", "Pricing/Models", "Incentives", "Security"], MEDIUM_BLUE),
        ("Batch\nProcessing", "11 programs", ["Daily/Weekly/Monthly", "CRM/DMS/GL Feeds", "Data Lake Extract", "Purge/Archive"], DARK_GRAY),
    ]

    for i, (name, count, items, color) in enumerate(domains):
        col = i % 4
        row = i // 4
        left = Inches(0.3) + Inches(col * 3.25)
        top = Inches(1.8) + Inches(row * 2.8)
        width = Inches(3.0)
        height = Inches(2.5)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Color header strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()

        add_text_box(slide, name, left + Inches(0.1), top + Inches(0.1),
                     Inches(1.8), Inches(0.55), font_size=13, color=color, bold=True)
        add_text_box(slide, count, left + Inches(2.0), top + Inches(0.15),
                     Inches(0.9), Inches(0.3), font_size=10, color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)

        for j, item in enumerate(items[:5]):
            add_text_box(slide, f"• {item}", left + Inches(0.1),
                         top + Inches(0.65) + Inches(j * 0.3),
                         width - Inches(0.2), Inches(0.3), font_size=10, color=DARK_GRAY)


def slide_05_architecture_decision(prs):
    """Architecture Decision - Why Modular Monolith."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Architecture Decision")
    add_subtitle_bar(slide, "Modular Monolith — the right pattern for deeply interdependent domains")

    # Three option cards
    options = [
        ("Option A: Microservices", [
            "8 separate services",
            "REST APIs between services",
            "Independent deployment",
            "PROBLEM: Deal completion",
            "touches 6+ domains — needs",
            "distributed transactions (saga)"
        ], RED, "Not recommended for this system"),
        ("Option B: Modular Monolith", [
            "Single app, clean module boundaries",
            "Shared DB, module-owned tables",
            "Internal method calls (no network)",
            "Single @Transactional boundary",
            "Can extract services later",
            "512MB RAM (vs 2GB+ for 8 services)"
        ], GREEN, "RECOMMENDED"),
        ("Option C: Event-Driven", [
            "Services via events (Kafka)",
            "Eventually consistent",
            "Loose coupling",
            "PROBLEM: Financial calcs need",
            "strong consistency — cannot",
            "tolerate eventual consistency"
        ], ORANGE, "Over-engineered for this use case"),
    ]

    for i, (title, items, color, verdict) in enumerate(options):
        left = Inches(0.3) + Inches(i * 4.3)
        card = add_card(slide, title, items, left, Inches(1.8),
                        width=Inches(4.0), height=Inches(3.2), accent=color)
        # Verdict badge
        add_text_box(slide, verdict, left + Inches(0.2), Inches(5.1),
                     Inches(3.6), Inches(0.4), font_size=12, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Key reasoning
    add_text_box(slide,
        "WHY: A single deal operation (sale completion) atomically updates Customer + Vehicle + Stock + "
        "Finance + Registration tables. In microservices, this requires complex saga patterns with compensation "
        "logic. In a modular monolith, it's a single @Transactional method — atomic, consistent, simple.",
        Inches(0.5), Inches(5.6), Inches(12), Inches(1.2),
        font_size=13, color=DARK_GRAY)

    # Prior project badge
    add_text_box(slide,
        "Note: Prior 3 projects used microservices (CICS-based, independent domains). "
        "This system's deep cross-domain dependencies make modular monolith the pragmatic choice.",
        Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
        font_size=11, color=MED_GRAY)


def slide_06_target_stack(prs):
    """Target Technology Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Target Technology Stack")
    add_subtitle_bar(slide, "IMS DC / COBOL / DB2 → Spring Boot / React / PostgreSQL")

    # Transformation table
    headers = ["Mainframe", "→", "Modern", "Notes"]
    rows = [
        ["IMS DC + MFS screens", "→", "React 18 + Tailwind CSS", "MFS → React pages with forms/tables"],
        ["IMS transaction codes", "→", "REST API endpoints", "ADMC → POST /api/admin/config"],
        ["COBOL programs", "→", "Spring Boot 3.3 (Java)", "Single modular monolith application"],
        ["DB2 z/OS (46+ tables)", "→", "PostgreSQL 16", "1:1 table migration, CHAR→VARCHAR"],
        ["IMS hierarchical (5 DBDs)", "→", "PostgreSQL (FK relationships)", "Segments → parent-child tables"],
        ["DL/I calls (GU/GN/ISRT)", "→", "JPA (find/save/delete)", "PCB status → JPA exceptions"],
        ["PSB/DBD definitions", "→", "JPA entities + Flyway DDL", "Segment sensitivity → entity fields"],
        ["COBOL subroutines (COM*)", "→", "Spring @Service beans", "Financial calcs ported to BigDecimal"],
        ["JCL batch + checkpoint", "→", "Spring @Scheduled + REST", "On-demand + scheduled execution"],
        ["EDI 214/856 (COMEDIL0)", "→", "Java EDI parser + webhook", "REST endpoint for carrier updates"],
        ["REXX/CLIST automation", "→", "Admin REST APIs + UI", "TSO utilities → web admin"],
        ["Copybooks (DCLGEN)", "→", "JPA entity classes", "Field layouts → class attributes"],
    ]
    col_widths = [Inches(2.8), Inches(0.4), Inches(2.8), Inches(4.5)]
    add_table_slide(slide, headers, rows, Inches(0.5), Inches(1.8), col_widths)


def slide_07_module_structure(prs):
    """Modular Monolith Structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Modular Monolith — Module Structure")
    add_subtitle_bar(slide, "Clean boundaries, shared transactions, extractable modules")

    modules = [
        ("common/", "Shared Infrastructure", ["JWT + Spring Security", "@Auditable AOP", "VIN Decoder", "Exception Handling"], DARK_GRAY),
        ("admin/", "Administration", ["Dealer Master", "Pricing & Models", "Tax Rates", "Incentives & Config"], MEDIUM_BLUE),
        ("customer/", "Customer Mgmt", ["CRUD + Search", "Credit Pre-Qual", "Lead Tracking", "Purchase History"], ACCENT_BLUE),
        ("sales/", "Sales Process", ["Deal Lifecycle", "Quote & Negotiate", "Trade-in & Incentive", "Approve & Complete"], GREEN),
        ("finance/", "Finance & FPL", ["Loan & Lease Calc", "Floor Plan Mgmt", "F&I Products", "Document Gen"], ORANGE),
        ("vehicle/", "Vehicle & Stock", ["Inventory Lifecycle", "Stock Position", "Shipment/Transit", "PDI & Location"], TEAL),
        ("registration/", "Reg & Warranty", ["DMV Submission", "Warranty Registration", "Recall Campaigns", "Claims Tracking"], PURPLE),
        ("batch/", "Batch & Integration", ["Daily/Weekly/Monthly", "CRM/DMS/GL Feeds", "EDI Parser", "Data Lake Extract"], RED),
    ]

    for i, (pkg, name, items, color) in enumerate(modules):
        col = i % 4
        row = i // 4
        left = Inches(0.3) + Inches(col * 3.25)
        top = Inches(1.8) + Inches(row * 2.7)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                      Inches(3.0), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Module name badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        left + Inches(0.1), top + Inches(0.1),
                                        Inches(1.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = tf.paragraphs[0].add_run()
        run.text = pkg
        run.font.size = Pt(9)
        run.font.color.rgb = WHITE
        run.font.bold = True

        add_text_box(slide, name, left + Inches(1.4), top + Inches(0.1),
                     Inches(1.5), Inches(0.35), font_size=12, color=color, bold=True)

        for j, item in enumerate(items):
            add_text_box(slide, f"• {item}", left + Inches(0.15),
                         top + Inches(0.55) + Inches(j * 0.38),
                         Inches(2.7), Inches(0.35), font_size=11, color=DARK_GRAY)


def slide_08_wave_plan(prs):
    """Wave Plan."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Wave Plan — 8 Waves")
    add_subtitle_bar(slide, "Domain-driven migration sequence, each wave delivers a complete vertical slice")

    waves = [
        ("Wave 0", "Foundation", "Auth + Common + DB Schema", DARK_GRAY, "1"),
        ("Wave 1", "Admin", "Dealer, Pricing, Tax, Models", MEDIUM_BLUE, "7"),
        ("Wave 2", "Customer", "CRUD, Credit, Leads, History", ACCENT_BLUE, "7"),
        ("Wave 3", "Sales", "Deal Lifecycle (Quote → Complete)", GREEN, "8"),
        ("Wave 4", "Finance", "Loan/Lease, Floor Plan, F&I", ORANGE, "12"),
        ("Wave 5", "Vehicle", "Inventory, Stock, Production, Logistics", TEAL, "24+"),
        ("Wave 6", "Registration", "Reg, Warranty, Recall", PURPLE, "11"),
        ("Wave 7", "Batch", "Scheduled Jobs, Integrations", RED, "11"),
    ]

    # Timeline arrow
    arrow_y = Inches(2.0)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), arrow_y,
                                    Inches(12.3), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_BLUE
    arrow.line.fill.background()

    for i, (wave, name, desc, color, progs) in enumerate(waves):
        left = Inches(0.3) + Inches(i * 1.6)
        top = Inches(2.7)

        # Wave circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.35),
                                         Inches(1.85), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = tf.paragraphs[0].add_run()
        run.text = str(i)
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True

        # Wave card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                      Inches(1.5), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Top accent
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(1.5), Inches(0.06))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()

        add_text_box(slide, wave, left + Inches(0.1), top + Inches(0.15),
                     Inches(1.3), Inches(0.3), font_size=10, color=MED_GRAY)
        add_text_box(slide, name, left + Inches(0.1), top + Inches(0.45),
                     Inches(1.3), Inches(0.35), font_size=14, color=color, bold=True)
        add_text_box(slide, desc, left + Inches(0.1), top + Inches(0.9),
                     Inches(1.3), Inches(1.2), font_size=10, color=DARK_GRAY)
        add_text_box(slide, f"{progs} programs", left + Inches(0.1), top + Inches(2.8),
                     Inches(1.3), Inches(0.3), font_size=10, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_text_box(slide, "Each wave delivers: Database schema + Backend module + REST APIs + React pages + Tests",
                 Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 font_size=12, color=MED_GRAY)


def slide_09_data_migration(prs):
    """Database Migration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Database Migration Strategy")
    add_subtitle_bar(slide, "Dual data stores (IMS hierarchical + DB2 relational) → Single PostgreSQL database")

    # IMS hierarchical mapping
    add_text_box(slide, "IMS Hierarchical → PostgreSQL Relational", Inches(0.5), Inches(1.7),
                 Inches(6), Inches(0.4), font_size=16, color=MEDIUM_BLUE, bold=True)

    ims_mappings = [
        "DBDAUTO1: VEHICLE (root) → vehicle + vehicle_option + vehicle_status_hist + lot_location",
        "DBDAUTO2: CUSTOMER (root) → customer + customer_lead + sales_deal + deal_line_item + trade_in",
        "DBDAUTO3: FINANCEAPP (root) → finance_app + finance_product + lease_terms + floor_plan_*",
        "DBDAUTO4: PRICEMASTER (root) → price_master + tax_rate + incentive_program",
        "DBDAUTO5: REPORTING → monthly_snapshot + stock_snapshot + audit_log",
    ]
    for i, m in enumerate(ims_mappings):
        add_text_box(slide, f"• {m}", Inches(0.5), Inches(2.2) + Inches(i * 0.35),
                     Inches(12), Inches(0.35), font_size=11, color=DARK_GRAY)

    # DB2 direct migration
    add_text_box(slide, "DB2 → PostgreSQL (Direct Migration)", Inches(0.5), Inches(4.2),
                 Inches(6), Inches(0.4), font_size=16, color=MEDIUM_BLUE, bold=True)

    headers = ["DB2 Type", "PostgreSQL", "Key Decision"]
    rows = [
        ["CHAR(n)", "VARCHAR(n)", "Avoids Hibernate CHAR padding bug (known from 3 prior projects)"],
        ["DECIMAL(p,s)", "NUMERIC(p,s)", "Financial precision preserved — BigDecimal in Java"],
        ["TIMESTAMP", "TIMESTAMP", "Direct mapping"],
        ["INTEGER", "INTEGER/BIGINT", "Direct mapping"],
    ]
    col_widths = [Inches(2), Inches(2), Inches(6.5)]
    add_table_slide(slide, headers, rows, Inches(0.5), Inches(4.8), col_widths)

    # Summary stat
    add_stat_box(slide, "Total PostgreSQL Tables", "~46",
                 Inches(9.5), Inches(1.8), width=Inches(3.3), height=Inches(1.2), color=TEAL)


def slide_10_enhancements(prs):
    """Enhancements Beyond Parity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Enhancements Beyond Mainframe Parity")
    add_subtitle_bar(slide, "12 features that demonstrate modernization adds tangible business value")

    enhancements = [
        ("Sales Pipeline Dashboard", "Kanban board: deals by stage with visual tracking", GREEN),
        ("Inventory Health Dashboard", "Aging heat map, stock turn, days-supply by model", TEAL),
        ("Finance Calculator Suite", "Side-by-side loan vs lease with what-if scenarios", ORANGE),
        ("Floor Plan Exposure View", "Total exposure by lender, curtailment alerts", ACCENT_BLUE),
        ("Audit Trail (AOP)", "@Auditable on all operations, searchable log", MEDIUM_BLUE),
        ("Excel/CSV Export", "Every list view exportable to XLSX/CSV", PURPLE),
        ("Real-time Deal Profitability", "Live gross calculation as terms change", GREEN),
        ("VIN Instant Decode", "Type VIN → immediate decode results", TEAL),
        ("On-demand Reports", "Run any batch report instantly", ORANGE),
        ("System Health Dashboard", "Batch job status, integration health", ACCENT_BLUE),
        ("User Management UI", "RBAC administration with role assignment", MEDIUM_BLUE),
        ("Paginated + Searchable Lists", "Server-side pagination, filter, sort", PURPLE),
    ]

    for i, (name, desc, color) in enumerate(enhancements):
        col = i % 3
        row = i // 3
        left = Inches(0.3) + Inches(col * 4.25)
        top = Inches(1.8) + Inches(row * 1.3)

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05),
                                        Inches(0.35), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(10)
        run.font.color.rgb = WHITE
        run.font.bold = True

        add_text_box(slide, name, left + Inches(0.45), top,
                     Inches(3.5), Inches(0.35), font_size=13, color=DARK_GRAY, bold=True)
        add_text_box(slide, desc, left + Inches(0.45), top + Inches(0.35),
                     Inches(3.5), Inches(0.35), font_size=10, color=MED_GRAY)


def slide_11_risk_register(prs):
    """Risk Register."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Risk Register & Mitigations")
    add_subtitle_bar(slide, "Key risks identified from codebase analysis and 3 prior modernization projects")

    headers = ["#", "Risk", "Severity", "Likelihood", "Mitigation"]
    rows = [
        ["1", "IMS hierarchical data flattening", "High", "Medium", "Map every DBD segment; validate with test data"],
        ["2", "Financial calc precision divergence", "High", "Low", "Port exact COBOL algorithms; BigDecimal; validate"],
        ["3", "Deal state machine edge cases", "High", "Medium", "State transition matrix; test every path"],
        ["4", "EDI parsing regression", "Medium", "Low", "Port COMEDIL0 logic; test with sample messages"],
        ["5", "Floor plan day-count variance", "Medium", "Medium", "Support all 3 bases (360/365/ACT); validate"],
        ["6", "Multi-jurisdiction tax complexity", "Medium", "Low", "Port COMTAXL0; test representative combos"],
        ["7", "8GB RAM constraint", "Medium", "High", "Monolith advantage: 1 JVM (~512MB)"],
        ["8", "Hibernate CHAR padding (known)", "Low", "High", "All VARCHAR from Day 1"],
    ]
    col_widths = [Inches(0.4), Inches(3.2), Inches(1.0), Inches(1.2), Inches(5.2)]
    add_table_slide(slide, headers, rows, Inches(0.5), Inches(1.8), col_widths)


def slide_12_testing(prs):
    """Testing Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Testing Strategy")
    add_subtitle_bar(slide, "~290 automated tests — highest count in the modernization series")

    # Test pyramid boxes
    levels = [
        ("Unit Tests (Backend)", "~130", "Service layer, calculators, VIN validation", "JUnit 5 + Mockito", GREEN),
        ("Integration Tests", "~70", "Controller endpoints + database", "@SpringBootTest", ACCENT_BLUE),
        ("Frontend Tests", "~60", "React components + pages", "Vitest + Testing Library", ORANGE),
        ("Financial Precision", "~30", "Loan/lease/interest/tax calculations", "Parameterized JUnit", PURPLE),
    ]

    for i, (name, count, scope, tool, color) in enumerate(levels):
        left = Inches(0.5)
        top = Inches(1.8) + Inches(i * 1.2)

        # Count badge
        add_stat_box(slide, "", count, left, top, width=Inches(1.3), height=Inches(0.9), color=color)

        add_text_box(slide, name, left + Inches(1.5), top,
                     Inches(3), Inches(0.4), font_size=14, color=color, bold=True)
        add_text_box(slide, f"{scope}  |  {tool}", left + Inches(1.5), top + Inches(0.4),
                     Inches(5), Inches(0.4), font_size=11, color=MED_GRAY)

    # Known patterns card
    add_card(slide, "Bug Prevention (from 3 prior projects)",
             ["UUID suffix for test data IDs",
              "VARCHAR for all columns (no CHAR)",
              "scanBasePackages includes common",
              "isString() for mutable assertions",
              "BigDecimal for all financial values"],
             Inches(8), Inches(1.8), width=Inches(4.8), height=Inches(3.0), accent=RED)

    # Total badge
    add_stat_box(slide, "Target Total Tests", "~290",
                 Inches(8), Inches(5.2), width=Inches(4.8), height=Inches(1.2), color=GREEN)


def slide_13_cross_project(prs):
    """Cross-Project Learning & Acceleration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI-Assisted Modernization — Cross-Project Acceleration")
    add_subtitle_bar(slide, "Project #4 benefits from accumulated learnings across 3 completed modernizations")

    # Prior projects summary
    headers = ["Project", "Source", "Target", "Programs", "SLOC", "Tests", "Bugs", "Timeline"]
    rows = [
        ["#1 CardDemo", "CICS/VSAM", "9 microservices", "48", "~25K", "~120", "~15", "~3 days"],
        ["#2 GenApp", "CICS/VSAM", "4 microservices", "31", "~11K", "153", "9", "~2 days"],
        ["#3 Portfolio", "CICS/DB2/VSAM", "5 microservices", "38", "~30K", "117", "8", "~4 days"],
        ["#4 AUTOSALES", "IMS DC/DB2/IMS", "Modular monolith", "103+", "~84.5K", "~290", "<10", "~5-7 days"],
    ]
    col_widths = [Inches(1.5), Inches(1.3), Inches(1.8), Inches(0.9), Inches(0.9), Inches(0.7), Inches(0.7), Inches(1)]
    add_table_slide(slide, headers, rows, Inches(0.8), Inches(1.7), col_widths)

    # Acceleration areas
    add_text_box(slide, "Where Cross-Project Memory Saves Effort (~50-60% reduction)",
                 Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
                 font_size=14, color=MEDIUM_BLUE, bold=True)

    saved = [
        ("Architecture Decisions", "~90%", "Proven stack, adapt for IMS specifics"),
        ("Wave Sequencing", "~80%", "Auth → Entities → Transactions → Reports → Batch"),
        ("Bug Prevention", "~60%", "Fix-forward applied from Day 1 (VARCHAR, UUID, scanBase)"),
        ("Enhancement Playbook", "~70%", "Dashboard + audit + export chosen instantly"),
    ]

    for i, (area, pct, desc) in enumerate(saved):
        left = Inches(0.5) + Inches(i * 3.1)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.8),
                                      Inches(2.9), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()
        add_text_box(slide, pct, left + Inches(0.15), Inches(4.9),
                     Inches(2.6), Inches(0.4), font_size=20, color=GREEN, bold=True)
        add_text_box(slide, area, left + Inches(0.15), Inches(5.3),
                     Inches(2.6), Inches(0.3), font_size=12, color=DARK_GRAY, bold=True)
        add_text_box(slide, desc, left + Inches(0.15), Inches(5.6),
                     Inches(2.6), Inches(0.5), font_size=10, color=MED_GRAY)

    # New effort note
    add_text_box(slide, "New: IMS DC patterns (MFS, DL/I, PSB/DBD, EDI) require full analysis — first non-CICS system",
                 Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
                 font_size=11, color=ORANGE)


def slide_14_deliverables(prs):
    """Deliverables Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Deliverables Summary")
    add_subtitle_bar(slide, "What you get: working application + comprehensive documentation + presentation assets")

    cards = [
        ("Working Application", [
            "Spring Boot modular monolith",
            "React SPA with 10+ pages",
            "PostgreSQL with 46 tables",
            "JWT authentication + RBAC",
            "~290 automated tests",
            "12 enhancements beyond parity"
        ], GREEN),
        ("Legacy Documentation", [
            "Codebase overview & inventory",
            "103+ program-level docs",
            "10 domain data flow docs",
            "Business rules extraction",
            "Cross-program synthesis",
            "Dependency mapping"
        ], ACCENT_BLUE),
        ("Modern App Documentation", [
            "API reference (all endpoints)",
            "Architecture overview",
            "Developer setup guide",
            "Database schema + mappings",
            "COBOL→Modern traceability",
            "Test inventory & strategy"
        ], ORANGE),
    ]

    for i, (title, items, color) in enumerate(cards):
        add_card(slide, title, items, Inches(0.3) + Inches(i * 4.3), Inches(1.8),
                 width=Inches(4.0), height=Inches(4.0), accent=color)

    # Total stats
    stats = [
        ("Total Files", "130+", ACCENT_BLUE),
        ("REST APIs", "80+", GREEN),
        ("React Pages", "20+", ORANGE),
        ("Documentation", "120+", PURPLE),
    ]
    for i, (label, value, color) in enumerate(stats):
        add_stat_box(slide, label, value,
                     Inches(0.3) + Inches(i * 3.25), Inches(6.2),
                     width=Inches(3.0), height=Inches(1.0), color=color)


def slide_15_call_to_action(prs):
    """Call to Action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(2.5), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(slide, "READY TO MODERNIZE", Inches(1), Inches(1.5),
                 Inches(11), Inches(0.6), font_size=18, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, "103+ COBOL Programs.\n10 Business Domains.\nOne Modular Monolith.",
                 Inches(1), Inches(2.8), Inches(11), Inches(1.8),
                 font_size=32, color=WHITE, bold=True)

    # Key points
    points = [
        "Proven approach — 4th modernization, 220+ COBOL programs modernized across 4 industries",
        "Right architecture — Modular monolith for deeply interdependent automotive sales domains",
        "Accelerated delivery — 50-60% faster from cross-project learnings and AI-assisted development",
        "Beyond parity — 12 enhancements that unlock business value inaccessible on the mainframe",
    ]
    for i, point in enumerate(points):
        add_text_box(slide, f"→  {point}", Inches(1), Inches(4.8) + Inches(i * 0.45),
                     Inches(11), Inches(0.4), font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

    add_text_box(slide, "Project #4  •  Automotive Industry  •  AI-Assisted Modernization Series",
                 Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                 font_size=12, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_title(prs)
    slide_02_exec_summary(prs)
    slide_03_current_state(prs)
    slide_04_business_domains(prs)
    slide_05_architecture_decision(prs)
    slide_06_target_stack(prs)
    slide_07_module_structure(prs)
    slide_08_wave_plan(prs)
    slide_09_data_migration(prs)
    slide_10_enhancements(prs)
    slide_11_risk_register(prs)
    slide_12_testing(prs)
    slide_13_cross_project(prs)
    slide_14_deliverables(prs)
    slide_15_call_to_action(prs)

    output_path = os.path.join(os.path.dirname(__file__),
                                "AUTOSALES_Modernization_Approval_Deck.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
