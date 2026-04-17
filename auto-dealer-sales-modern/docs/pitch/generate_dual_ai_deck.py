#!/usr/bin/env python3
"""
AUTOSALES — Dual-AI Pitch Deck Generator
Produces AUTOSALES_Dual_AI_Pitch.pptx — a standalone visual pitch for the
AI Assistant + AI Agent story (shipped 2026-04-14).

Visual language
---------------
- Blue family  = AI Assistant (fast, free-tier, single-step)
- Violet family = AI Agent (premium, Claude, multi-step)
- Dark navy header bar on every content slide
- Rounded cards, colored pills, shape-based diagrams (no ASCII art)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# ── Palette ───────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)
DEEP_NAVY   = RGBColor(0x07, 0x12, 0x20)
STEEL       = RGBColor(0x1B, 0x4F, 0x72)

BLUE        = RGBColor(0x21, 0x96, 0xF3)
BLUE_DARK   = RGBColor(0x0D, 0x47, 0xA1)
BLUE_SOFT   = RGBColor(0xE3, 0xF2, 0xFD)

VIOLET      = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_DARK = RGBColor(0x4C, 0x1D, 0x95)
VIOLET_SOFT = RGBColor(0xEE, 0xE6, 0xFF)

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xFA, 0xFA, 0xFD)
DARK_GRAY   = RGBColor(0x2B, 0x2B, 0x2B)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
SOFT_GRAY   = RGBColor(0xE5, 0xE7, 0xEB)

GREEN       = RGBColor(0x22, 0xC5, 0x5E)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)
TEAL        = RGBColor(0x06, 0xB6, 0xD4)
PINK        = RGBColor(0xEC, 0x48, 0x99)

SW = Inches(13.333)
SH = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────────
def _i(v):
    # Coerce an EMU-ish value to int; PowerPoint's OOXML requires integer EMUs
    # and float divisions (e.g. SW/2, h/2) produce floats that corrupt the file.
    return int(round(v)) if v is not None else v


def _run(p, text, size=14, bold=False, color=DARK_GRAY, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def textbox(slide, text, left, top, width, height, size=14, bold=False,
            color=DARK_GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    left, top, width, height = _i(left), _i(top), _i(width), _i(height)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size=size, bold=bold, color=color, italic=italic)
    return tf


def rect(slide, left, top, width, height, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    left, top, width, height = _i(left), _i(top), _i(width), _i(height)
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def pill(slide, text, left, top, width, height, fill, fg=WHITE, size=12, bold=True):
    s = rect(slide, left, top, width, height, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = s.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, text, size=size, bold=bold, color=fg)
    return s


def header(slide, title, subtitle=None):
    # Main bar
    rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    textbox(slide, title, Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.55),
            size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Accent stripe: half blue, half violet (the whole deck theme)
    rect(slide, 0, Inches(0.9), SW / 2, Inches(0.08), BLUE)
    rect(slide, SW / 2, Inches(0.9), SW / 2, Inches(0.08), VIOLET)
    if subtitle:
        textbox(slide, subtitle, Inches(0.5), Inches(1.05), SW - Inches(1), Inches(0.4),
                size=14, color=MED_GRAY, italic=True)


def footer(slide, page_num, total=9):
    textbox(slide, "AUTOSALES · Dual-AI Pitch · 2026-04", Inches(0.3), SH - Inches(0.35),
            Inches(5), Inches(0.3), size=9, color=MED_GRAY)
    textbox(slide, f"{page_num} / {total}", SW - Inches(1), SH - Inches(0.35),
            Inches(0.7), Inches(0.3), size=9, color=MED_GRAY, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, left, top, width, height, size=13, color=DARK_GRAY,
                 bullet_color=None, line_spacing=1.25):
    left, top, width, height = _i(left), _i(top), _i(width), _i(height)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        _run(p, "● ", size=size, bold=True, color=bullet_color or color)
        _run(p, item, size=size, color=color)
    return tf


def arrow(slide, x1, y1, x2, y2, color=MED_GRAY, weight=Pt(2)):
    # Plain line segment; callers use textbox glyphs/labels for direction cues.
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = weight
    return c


# ── Slides ────────────────────────────────────────────────────────
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-bleed gradient-style: navy top 70%, blue+violet bottom strip
    rect(s, 0, 0, SW, Inches(5.5), DEEP_NAVY)
    rect(s, 0, Inches(5.5), SW / 2, Inches(2.0), BLUE_DARK)
    rect(s, SW / 2, Inches(5.5), SW / 2, Inches(2.0), VIOLET_DARK)

    # Eyebrow
    textbox(s, "AUTOSALES MODERNIZATION", Inches(0.8), Inches(1.1),
            Inches(10), Inches(0.4), size=14, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

    # Title
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(12), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _run(p, "Two AI Surfaces,", size=60, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    _run(p2, "One Dealership.", size=60, bold=True, color=BLUE)

    # Subtitle
    textbox(s, "How AUTOSALES moved from green-screen transactions to a system that thinks alongside the dealer.",
            Inches(0.8), Inches(4.4), Inches(12), Inches(0.7),
            size=18, color=SOFT_GRAY, italic=True)

    # Bottom strip pills
    pill(s, "AI ASSISTANT  ·  /api/chat  ·  4 free-tier LLMs",
         Inches(0.8), Inches(6.1), Inches(5.8), Inches(0.55), BLUE, WHITE, 14, True)
    pill(s, "AI AGENT  ·  /api/agent  ·  Claude Sonnet 4.6",
         Inches(6.75), Inches(6.1), Inches(5.8), Inches(0.55), VIOLET, WHITE, 14, True)

    textbox(s, "Shipped 2026-04-14  ·  OpenClaw skills architecture  ·  Level 4 agentic",
            Inches(0.8), Inches(6.85), Inches(12), Inches(0.4),
            size=12, color=SOFT_GRAY, italic=True)


def slide_ladder(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "The Evolution Ladder",
           "Modernization isn't just replacing screens — it's giving the dealer a system that reasons.")

    # 4 columns, each a vertical card climbing in color intensity
    cols = [
        ("IMS DC\n(1990s)",       "Green screen, PF keys",      "Clerk knows every screen & txn code", "System executes the exact txn",        SOFT_GRAY, DARK_GRAY),
        ("React UI\n(2026)",      "Forms, pages, filters",      "Clerk clicks through workflows",      "System shows data, validates input",    BLUE_SOFT, STEEL),
        ("AI Assistant\n(2026)",  "Chat box",                   "Asks ONE question",                   "Fetches ONE answer",                    BLUE, WHITE),
        ("AI Agent\n(2026)",      "Chat box",                   "States a GOAL",                       "Plans, chains tools, reasons, recommends", VIOLET, WHITE),
    ]
    left0 = Inches(0.5); top = Inches(1.7); w = Inches(3.05); h = Inches(4.3); gap = Inches(0.12)
    for i, (era, mode, human, sys, bg, fg) in enumerate(cols):
        x = left0 + (w + gap) * i
        card = rect(s, x, top, w, h, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # header band inside card
        band = rect(s, x, top, w, Inches(1.1), fg if bg in (BLUE, VIOLET) else STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, era, x, top + Inches(0.1), w, Inches(0.95),
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # body
        textbox(s, "INTERACTION", x + Inches(0.2), top + Inches(1.3), w - Inches(0.4), Inches(0.3),
                size=9, bold=True, color=MED_GRAY)
        textbox(s, mode, x + Inches(0.2), top + Inches(1.55), w - Inches(0.4), Inches(0.6),
                size=12, color=fg)

        textbox(s, "HUMAN", x + Inches(0.2), top + Inches(2.25), w - Inches(0.4), Inches(0.3),
                size=9, bold=True, color=MED_GRAY)
        textbox(s, human, x + Inches(0.2), top + Inches(2.5), w - Inches(0.4), Inches(0.7),
                size=12, color=fg, bold=True)

        textbox(s, "SYSTEM", x + Inches(0.2), top + Inches(3.2), w - Inches(0.4), Inches(0.3),
                size=9, bold=True, color=MED_GRAY)
        textbox(s, sys, x + Inches(0.2), top + Inches(3.45), w - Inches(0.4), Inches(0.75),
                size=12, color=fg)

    # Punchline banner
    rect(s, Inches(0.5), Inches(6.25), SW - Inches(1), Inches(0.7), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "The Agent is the first time the software does the thinking the clerk used to do.",
            Inches(0.5), Inches(6.25), SW - Inches(1), Inches(0.7),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)


def _surface_card(slide, x, y, w, h, accent, accent_soft, accent_dark,
                  title, tagline, endpoint, llm, bullets, sample_q, sample_a, latency):
    # Outer card
    rect(slide, x, y, w, h, WHITE, line=SOFT_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # Accent top band
    rect(slide, x, y, w, Inches(1.1), accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # White overlay so only the top stays colored (hack: big rect below band)
    rect(slide, x, y + Inches(0.55), w, Inches(0.55), accent)

    textbox(slide, title, x + Inches(0.3), y + Inches(0.12), w - Inches(0.6), Inches(0.5),
            size=22, bold=True, color=WHITE)
    textbox(slide, tagline, x + Inches(0.3), y + Inches(0.62), w - Inches(0.6), Inches(0.4),
            size=13, color=WHITE, italic=True)

    # Meta chips
    pill(slide, endpoint, x + Inches(0.3), y + Inches(1.3), Inches(2.1), Inches(0.35),
         accent_soft, accent_dark, 10, True)
    pill(slide, llm, x + Inches(2.5), y + Inches(1.3), Inches(3.2), Inches(0.35),
         accent_soft, accent_dark, 10, True)
    pill(slide, latency, x + w - Inches(1.5), y + Inches(1.3), Inches(1.2), Inches(0.35),
         accent, WHITE, 10, True)

    # Bullets
    bullet_block(slide, bullets, x + Inches(0.3), y + Inches(1.85), w - Inches(0.6), Inches(1.7),
                 size=12, color=DARK_GRAY, bullet_color=accent)

    # Sample Q/A block
    qy = y + Inches(3.7)
    rect(slide, x + Inches(0.3), qy, w - Inches(0.6), Inches(0.7), accent_soft,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(slide, "ASK", x + Inches(0.45), qy + Inches(0.08), Inches(0.6), Inches(0.25),
            size=8, bold=True, color=accent_dark)
    textbox(slide, sample_q, x + Inches(0.45), qy + Inches(0.3), w - Inches(0.9), Inches(0.4),
            size=12, color=DARK_GRAY, italic=True)

    ay = qy + Inches(0.85)
    rect(slide, x + Inches(0.3), ay, w - Inches(0.6), Inches(0.85), WHITE,
         line=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(slide, "ANSWER", x + Inches(0.45), ay + Inches(0.08), Inches(0.9), Inches(0.25),
            size=8, bold=True, color=accent_dark)
    textbox(slide, sample_a, x + Inches(0.45), ay + Inches(0.3), w - Inches(0.9), Inches(0.55),
            size=11, color=DARK_GRAY)


def slide_assistant(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Meet the Assistant", "Fast, cheap, single-step. Google for your dealership.")

    _surface_card(
        s, Inches(2.5), Inches(1.7), Inches(8.3), Inches(5.55),
        accent=BLUE, accent_soft=BLUE_SOFT, accent_dark=BLUE_DARK,
        title="AI ASSISTANT",
        tagline="\"Google for your dealership.\"",
        endpoint="POST /api/chat",
        llm="Groq · Gemini · Together · Mistral",
        bullets=[
            "Single LLM call, user picks the provider",
            "Free-tier economics — zero marginal cost per question",
            "Use when you know what you want, just need it faster than clicking",
            "Blue pill in header with provider dropdown",
        ],
        sample_q="What's customer 42's phone number?",
        sample_a="Customer 42 — John Doe — (555) 123-4567. Last contact 2026-04-10.",
        latency="~1–2 s",
    )
    footer(s, 3)


def slide_agent(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Meet the Agent", "Deliberate, multi-step, skills-based. A junior analyst on staff.")

    _surface_card(
        s, Inches(2.5), Inches(1.7), Inches(8.3), Inches(5.55),
        accent=VIOLET, accent_soft=VIOLET_SOFT, accent_dark=VIOLET_DARK,
        title="AI AGENT",
        tagline="\"A junior analyst on staff.\"",
        endpoint="POST /api/agent",
        llm="Claude Sonnet 4.6  via  OpenClaw gateway",
        bullets=[
            "Plans before acting — states the steps, then executes",
            "Chains 4–8 tool calls autonomously, no turn-by-turn prompting",
            "Applies domain rules (aging, APR bands, approval gates, SSN redaction)",
            "Surfaces non-obvious insights — referrals, stalled deals, risks",
        ],
        sample_q="Is deal 17 healthy?",
        sample_a="Stalled at UW for 9 days · credit check stale (28d) · FIN-APP missing APR band · 3 recommended actions.",
        latency="20–40 s",
    )
    footer(s, 4)


def slide_axes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Four Differentiation Axes", "Why the Agent is a new product category, not just a fancier chatbot.")

    axes = [
        ("AI",  "Cognitive Load",
         "Assistant answers YOUR question.",
         "Agent decides WHICH questions matter.",
         BLUE, VIOLET),
        (">>",  "Tool-Call Depth",
         "1 API call.",
         "4-8 chained calls with domain rules between steps.",
         BLUE, VIOLET),
        ("$",   "Cost & Latency",
         "Free - ~2 s.",
         "~$0.05 - 20-40 s - value gap justifies trade.",
         BLUE, VIOLET),
        ("/\\", "Architecture",
         "Single LLM call, one prompt.",
         "Skills-based: SKILL.md + 28 tools + 5 recipes + OpenClaw.",
         BLUE, VIOLET),
    ]
    # 2x2 grid
    left0 = Inches(0.6); top0 = Inches(1.7); w = Inches(6.0); h = Inches(2.6); gx = Inches(0.15); gy = Inches(0.2)
    for i, (icon, name, a_line, g_line, c_blue, c_violet) in enumerate(axes):
        r, c = divmod(i, 2)
        x = left0 + (w + gx) * c
        y = top0 + (h + gy) * r
        rect(s, x, y, w, h, WHITE, line=SOFT_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Left index ribbon
        rect(s, x, y, Inches(0.55), h, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, f"0{i+1}", x, y, Inches(0.55), h,
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title row
        textbox(s, f"{icon}  {name}", x + Inches(0.7), y + Inches(0.15), w - Inches(0.9), Inches(0.5),
                size=18, bold=True, color=NAVY)

        # Split rows
        rect(s, x + Inches(0.7), y + Inches(0.85), w - Inches(0.9), Inches(0.72), BLUE_SOFT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        pill(s, "ASSISTANT", x + Inches(0.85), y + Inches(0.98), Inches(1.1), Inches(0.3),
             BLUE, WHITE, 9, True)
        textbox(s, a_line, x + Inches(2.05), y + Inches(0.95), w - Inches(2.3), Inches(0.55),
                size=12, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

        rect(s, x + Inches(0.7), y + Inches(1.68), w - Inches(0.9), Inches(0.78), VIOLET_SOFT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        pill(s, "AGENT", x + Inches(0.85), y + Inches(1.83), Inches(1.1), Inches(0.3),
             VIOLET, WHITE, 9, True)
        textbox(s, g_line, x + Inches(2.05), y + Inches(1.78), w - Inches(2.3), Inches(0.65),
                size=12, color=DARK_GRAY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 5)


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Same Question, Two Answers", "Prompt both surfaces with: \"Give me a morning briefing.\"")

    # Prompt banner
    rect(s, Inches(0.5), Inches(1.75), SW - Inches(1), Inches(0.7), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "PROMPT:  \"Give me a morning briefing.\"",
            Inches(0.5), Inches(1.75), SW - Inches(1), Inches(0.7),
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Two side-by-side answer cards
    top = Inches(2.7); h = Inches(4.3); w = Inches(6.1)
    # Assistant
    x1 = Inches(0.5)
    rect(s, x1, top, w, h, WHITE, line=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x1, top, w, Inches(0.75), BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x1, top + Inches(0.4), w, Inches(0.35), BLUE)
    textbox(s, "ASSISTANT  -  ~2 s  -  1 call",
            x1 + Inches(0.25), top + Inches(0.1), w - Inches(0.5), Inches(0.55),
            size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    textbox(s, "\"Here are today's deals:\"",
            x1 + Inches(0.3), top + Inches(0.95), w - Inches(0.6), Inches(0.4),
            size=13, italic=True, color=MED_GRAY)
    # Fake list rows
    rows = [
        "• Deal #14 — Johnson — WS",
        "• Deal #15 — Garcia — CA",
        "• Deal #16 — Patel — UW",
        "• Deal #17 — Johnson — UW",
        "• Deal #18 — Chen — DL",
        "• Deal #19 — O'Neil — WS",
        "… 11 more rows",
    ]
    bullet_block(s, rows, x1 + Inches(0.3), top + Inches(1.45), w - Inches(0.6), Inches(2.6),
                 size=12, color=DARK_GRAY, bullet_color=BLUE)
    # Footer label
    rect(s, x1, top + h - Inches(0.5), w, Inches(0.5), BLUE_SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "Raw data. Reader must interpret.",
            x1, top + h - Inches(0.5), w, Inches(0.5),
            size=12, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Agent
    x2 = Inches(6.75)
    rect(s, x2, top, w, h, WHITE, line=VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x2, top, w, Inches(0.75), VIOLET, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x2, top + Inches(0.4), w, Inches(0.35), VIOLET)
    textbox(s, "AGENT  -  ~30 s  -  6 chained calls + rules",
            x2 + Inches(0.25), top + Inches(0.1), w - Inches(0.5), Inches(0.55),
            size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    textbox(s, "\"Good morning. 3 things need you today:\"",
            x2 + Inches(0.3), top + Inches(0.95), w - Inches(0.6), Inches(0.4),
            size=13, italic=True, color=MED_GRAY)

    insights = [
        ("1", RED,    "Deal #17 stalled at UW for 9 days. Credit check stale; customer Johnson last contacted Tuesday."),
        ("2", AMBER,  "VIN 1HG... on lot 72 days - 12 days over floor-plan exposure. Suggest transfer to DLR02."),
        ("3", GREEN,  "Warranty claim #44 crosses 14-day escalation threshold tomorrow - act today."),
    ]
    iy = top + Inches(1.45)
    for ic, col, txt in insights:
        rect(s, x2 + Inches(0.3), iy, w - Inches(0.6), Inches(0.75), VIOLET_SOFT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x2 + Inches(0.3), iy, Inches(0.55), Inches(0.75), col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, ic, x2 + Inches(0.3), iy, Inches(0.55), Inches(0.75),
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, txt, x2 + Inches(0.95), iy + Inches(0.05), w - Inches(1.25), Inches(0.7),
                size=11, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
        iy += Inches(0.85)

    rect(s, x2, top + h - Inches(0.5), w, Inches(0.5), VIOLET_SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "Prioritized decisions. Human acts, not interprets.",
            x2, top + h - Inches(0.5), w, Inches(0.5),
            size=12, bold=True, color=VIOLET_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 6)


def slide_recipes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Five Workflow Recipes",
           "Shipped in openclaw/skills/autosales-api/SKILL.full.md — each encodes a dealer's domain rules.")

    recipes = [
        ("HC",  "Deal Health Check",       "State machine + SLA aging + missing docs",              RED),
        ("C360","Customer 360",            "Referral graph + stale-data flags + loyalty signals",   BLUE),
        ("FN",  "Lead-to-Deal Funnel",     "Conversion diagnostics, stalled stages",                 TEAL),
        ("INV", "Inventory Aging Triage",  "60-day rule + floor-plan exposure + transfer hints",     AMBER),
        ("MB",  "Morning Briefing",        "Composite rollup across all four above",                 VIOLET),
    ]
    top = Inches(1.7); h = Inches(1.0); gap = Inches(0.12)
    for i, (icon, name, desc, col) in enumerate(recipes):
        y = top + (h + gap) * i
        rect(s, Inches(0.8), y, SW - Inches(1.6), h, WHITE, line=SOFT_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Accent bar
        rect(s, Inches(0.8), y, Inches(0.25), h, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Icon tile
        rect(s, Inches(1.2), y + Inches(0.15), Inches(0.7), Inches(0.7), col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, icon, Inches(1.2), y + Inches(0.15), Inches(0.7), Inches(0.7),
                size=22, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Index
        textbox(s, f"0{i+1}", Inches(2.05), y + Inches(0.08), Inches(0.5), Inches(0.3),
                size=10, bold=True, color=MED_GRAY)
        # Name
        textbox(s, name, Inches(2.05), y + Inches(0.28), Inches(4.0), Inches(0.45),
                size=18, bold=True, color=NAVY)
        # Desc
        textbox(s, desc, Inches(6.2), y + Inches(0.3), Inches(5.2), Inches(0.45),
                size=13, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
        # CTA pill
        pill(s, "SKILL.full.md", SW - Inches(2.5), y + Inches(0.3),
             Inches(1.6), Inches(0.4), col, WHITE, 10, True)
    footer(s, 7)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Under the Hood", "React → Spring Boot → OpenClaw gateway → Claude Sonnet 4.6, with skill-based progressive disclosure.")

    # Row 1: pipeline boxes
    y = Inches(2.0); h = Inches(1.4)
    boxes = [
        ("React UI",            "AgentWidget\n(violet pill)",       BLUE_DARK,   WHITE),
        ("Spring Boot",         "AgentController\nAgentService",    STEEL,       WHITE),
        ("OpenClaw Gateway",    ":18789\nrouter + plugin",          TEAL,        WHITE),
        ("Claude Sonnet 4.6",   "Anthropic plugin\nLLM reasoning",  VIOLET,      WHITE),
    ]
    w = Inches(2.7); gap = Inches(0.45)
    total = w * len(boxes) + gap * (len(boxes) - 1)
    x0 = (SW - total) / 2
    centers = []
    for i, (title, sub, bg, fg) in enumerate(boxes):
        x = x0 + (w + gap) * i
        rect(s, x, y, w, h, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, title, x, y + Inches(0.15), w, Inches(0.45),
                size=15, bold=True, color=fg, align=PP_ALIGN.CENTER)
        textbox(s, sub, x, y + Inches(0.6), w, Inches(0.75),
                size=11, color=fg, align=PP_ALIGN.CENTER)
        centers.append((x + w, y + h / 2, x, y + h / 2))

    # Arrows between boxes
    for i in range(len(boxes) - 1):
        right_x, mid_y, _, _ = centers[i]
        _, _, next_left, _ = centers[i + 1]
        arrow(s, right_x, mid_y, next_left, mid_y, color=MED_GRAY, weight=Pt(2.5))
        # Label
        lbl = ["POST /api/agent", "POST /v1/chat/completions", "Anthropic API"][i]
        textbox(s, lbl, right_x, mid_y - Inches(0.42), next_left - right_x, Inches(0.3),
                size=9, color=MED_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Row 2: Skill panel fed into gateway
    sy = Inches(4.1); sh_ = Inches(1.9)
    # Skill card (centered under gateway)
    gx = centers[2][2]  # gateway left
    skill_w = Inches(6.0)
    skill_x = (SW - skill_w) / 2
    rect(s, skill_x, sy, skill_w, sh_, OFF_WHITE, line=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, skill_x, sy, skill_w, Inches(0.5), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "Skill Package  -  openclaw/skills/autosales-api/",
            skill_x, sy, skill_w, Inches(0.5),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Two sub-cards: SKILL.md / SKILL.full.md
    sub_w = Inches(2.75); sub_h = Inches(1.15); sub_y = sy + Inches(0.6)
    rect(s, skill_x + Inches(0.15), sub_y, sub_w, sub_h, WHITE, line=SOFT_GRAY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "SKILL.md", skill_x + Inches(0.25), sub_y + Inches(0.1), sub_w, Inches(0.35),
            size=12, bold=True, color=TEAL)
    textbox(s, "Lean trigger (~30 lines).\nLoaded always.",
            skill_x + Inches(0.25), sub_y + Inches(0.45), sub_w, Inches(0.65),
            size=10, color=DARK_GRAY)

    rect(s, skill_x + sub_w + Inches(0.3), sub_y, sub_w, sub_h, WHITE, line=SOFT_GRAY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "SKILL.full.md", skill_x + sub_w + Inches(0.4), sub_y + Inches(0.1), sub_w, Inches(0.35),
            size=12, bold=True, color=VIOLET)
    textbox(s, "Payload (~230 lines). 28 tools,\n5 recipes, domain rules.",
            skill_x + sub_w + Inches(0.4), sub_y + Inches(0.45), sub_w, Inches(0.65),
            size=10, color=DARK_GRAY)

    # Arrow from skill card up to gateway
    gw_x = centers[2][2] + w / 2  # gateway center x
    gw_y = y + h  # gateway bottom
    arrow(s, gw_x, gw_y, skill_x + skill_w / 2, sy, color=TEAL, weight=Pt(2))
    textbox(s, "progressive disclosure", gw_x - Inches(1.0), (gw_y + sy) / 2 - Inches(0.15),
            Inches(2.0), Inches(0.3), size=9, color=TEAL, italic=True, align=PP_ALIGN.CENTER)

    # AUTOSALES API box (bottom right) — tools
    api_x = Inches(9.3); api_y = Inches(6.2); api_w = Inches(3.5); api_h = Inches(0.8)
    rect(s, api_x, api_y, api_w, api_h, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, "AUTOSALES REST API  -  28 tools",
            api_x, api_y, api_w, api_h,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Connect skill to API
    arrow(s, skill_x + skill_w / 2, sy + sh_, api_x + api_w / 2, api_y,
          color=VIOLET, weight=Pt(2))

    # Auth note
    textbox(s, "Auth:  JWT at /api/agent  -  X-API-Key gateway->app  -  Anthropic key in .env",
            Inches(0.5), api_y + Inches(0.2), Inches(8.5), Inches(0.3),
            size=10, color=MED_GRAY, italic=True)
    footer(s, 8)


def slide_whats_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "What's Next", "Polish, depth, and operational guardrails for the agent surface.")

    items = [
        ("UI",  "Tool-Call Visibility",        "SSE streaming so the agent's plan shows live as it runs.",            BLUE),
        ("MEM", "Server-Side Memory",          "Per-userId history in Postgres. Enables longer multi-turn.",          VIOLET),
        ("CT",  "Composite Tools",             "get_customer_360(id) as one orchestrated call - fewer round-trips.",   TEAL),
        ("RCP", "More Recipes",                "Finance deal review - Inventory rebalance - Recall impact report.",    AMBER),
        ("WEB", "External web_fetch",          "Pull NHTSA recalls, KBB pricing, VIN-decode enhancements.",            GREEN),
        ("$$",  "Cost Guardrails",             "Monthly cap, per-user quotas, budget alerts.",                          RED),
    ]
    # 3x2 grid of tiles
    left0 = Inches(0.6); top0 = Inches(1.7); w = Inches(4.0); h = Inches(2.4); gx = Inches(0.2); gy = Inches(0.25)
    for i, (icon, name, desc, col) in enumerate(items):
        r, c = divmod(i, 3)
        x = left0 + (w + gx) * c
        y = top0 + (h + gy) * r
        rect(s, x, y, w, h, WHITE, line=SOFT_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Icon tile
        rect(s, x + Inches(0.3), y + Inches(0.3), Inches(1.0), Inches(1.0), col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, icon, x + Inches(0.3), y + Inches(0.3), Inches(1.0), Inches(1.0),
                size=28, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, name, x + Inches(1.45), y + Inches(0.35), w - Inches(1.65), Inches(0.5),
                size=16, bold=True, color=NAVY)
        textbox(s, desc, x + Inches(0.3), y + Inches(1.45), w - Inches(0.6), Inches(0.9),
                size=12, color=DARK_GRAY)
    footer(s, 9)


# ── Main ──────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_title(prs)         # 1
    slide_ladder(prs)        # 2
    slide_assistant(prs)     # 3
    slide_agent(prs)         # 4
    slide_axes(prs)          # 5
    slide_demo(prs)          # 6
    slide_recipes(prs)       # 7
    slide_architecture(prs)  # 8
    slide_whats_next(prs)    # 9

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "AUTOSALES_Dual_AI_Pitch.pptx")
    prs.save(out)
    print(f"OK -> {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
